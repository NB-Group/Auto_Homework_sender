import os
import json
import requests
import sys
from pathlib import Path
from pptx import Presentation
from tkinter import filedialog
import threading
import time

def get_app_data_path():
    """获取应用数据存储路径（优先使用用户AppData\Auto_Homework）。"""
    try:
        base = os.getenv('APPDATA') or os.getenv('LOCALAPPDATA') or os.path.expanduser('~')
        app_dir = os.path.join(base, 'Auto_Homework')
        os.makedirs(app_dir, exist_ok=True)
        return app_dir
    except Exception:
        # 回退：打包后使用可执行文件所在目录；否则使用脚本目录
        try:
            if getattr(sys, 'frozen', False) or hasattr(sys, '_MEIPASS'):
                return os.path.dirname(sys.executable)
        except Exception:
            pass
        return os.path.abspath(os.path.dirname(__file__))

class HomeworkAPI:
    def __init__(self):
        self.config_file = os.path.join(get_app_data_path(), "config.json")
        self.config = self.load_config()
        
    def load_config(self):
        """加载配置文件"""
        # 优先从 AppData 读取；若不存在，尝试从旧位置迁移
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass

        # 尝试从旧位置读取并迁移（工作目录或脚本目录）
        try:
            legacy_candidates = [
                os.path.join(os.path.abspath('.'), 'config.json'),
                os.path.join(os.path.abspath(os.path.dirname(__file__)), 'config.json')
            ]
            for legacy in legacy_candidates:
                if os.path.exists(legacy):
                    with open(legacy, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                    # 写入到新位置
                    try:
                        with open(self.config_file, 'w', encoding='utf-8') as wf:
                            json.dump(data, wf, ensure_ascii=False, indent=2)
                    except Exception:
                        pass
                    return data
        except Exception:
            pass

        return {
                "access_token": "75a4cb7b5e71dbc785977184dcbaf0e11b7c355a5748406d158b464e85d62637",
                "auto_send_time": "17:00",
                "auto_send_enabled": False,
                "ppt_file_path": "",
                "theme": "dark",
                "weekday_send_time": "17:00",
                "friday_send_time": "15:00"
            }
    
    def save_config(self, config):
        """保存配置"""
        self.config.update(config)
        # 确保目录存在
        try:
            os.makedirs(os.path.dirname(self.config_file), exist_ok=True)
        except Exception:
            pass
        with open(self.config_file, 'w', encoding='utf-8') as f:
            json.dump(self.config, f, ensure_ascii=False, indent=2)
        return {"success": True}
    
    def select_ppt_file(self):
        """选择PPT文件"""
        try:
            # 如果配置中已有文件路径，使用该路径作为初始目录
            initial_dir = ""
            if self.config.get("ppt_file_path"):
                initial_dir = os.path.dirname(self.config["ppt_file_path"])
            
            file_path = filedialog.askopenfilename(
                title="选择PPT文件",
                filetypes=[("PowerPoint files", "*.pptx *.ppt")],
                initialdir=initial_dir
            )
            
            if file_path:
                # 保存选择的文件路径到配置
                self.config["ppt_file_path"] = file_path
                self.save_config(self.config)
            
            return {"success": True, "file_path": file_path}
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def parse_ppt_to_markdown(self, file_path):
        """将PPT转换为Markdown格式"""
        try:
            presentation = Presentation(file_path)
            markdown_content = "# 作业内容\n\n"
            
            # 只处理最后一页
            if presentation.slides:
                last_slide = presentation.slides[-1]
                
                for shape in last_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 每行一个无序列表项（不额外添加空白行）
                        normalized = shape.text.replace('\r', '\n')
                        lines = [ln.strip() for ln in normalized.split('\n') if ln.strip()]
                        for ln in lines:
                            markdown_content += f"- {ln}\n"

            # 去除结尾多余空白行，仅保留必要结尾换行
            markdown_content = markdown_content.rstrip() + "\n"
            
            return {"success": True, "content": markdown_content}
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def send_to_dingtalk(self, content, access_token=None):
        """发送消息到钉钉"""
        try:
            token = access_token or self.config.get("access_token")
            if not token:
                return {"success": False, "error": "未设置ACCESS_TOKEN"}
            
            # 确保消息包含"作业"字段
            if "作业" not in content:
                content = f"【作业通知】\n\n{content}"
            
            url = f"https://oapi.dingtalk.com/robot/send?access_token={token}"
            
            payload = {
                "msgtype": "markdown",
                "markdown": {
                    "title": "作业通知",
                    "text": content
                }
            }
            
            response = requests.post(url, json=payload, timeout=10)
            result = response.json()
            
            if result.get("errcode") == 0:
                return {"success": True, "message": "发送成功"}
            else:
                return {"success": False, "error": f"发送失败: {result.get('errmsg')}"}
                
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def preview_homework(self, file_path):
        """预览作业内容"""
        result = self.parse_ppt_to_markdown(file_path)
        if result["success"]:
            return {
                "success": True,
                "content": result["content"],
                "file_name": os.path.basename(file_path)
            }
        return result
    
    def send_homework(self, file_path):
        """发送作业"""
        parse_result = self.parse_ppt_to_markdown(file_path)
        if not parse_result["success"]:
            return parse_result
        
        send_result = self.send_to_dingtalk(parse_result["content"])
        return send_result
    
    def auto_send_homework(self):
        """自动发送作业（使用配置中的PPT文件路径）"""
        try:
            ppt_path = self.config.get("ppt_file_path")
            if not ppt_path or not os.path.exists(ppt_path):
                return {"success": False, "error": "未设置PPT文件路径或文件不存在"}
            
            return self.send_homework(ppt_path)
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def get_config(self):
        """获取当前配置"""
        return self.config
